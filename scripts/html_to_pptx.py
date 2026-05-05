from __future__ import annotations

import argparse
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BLUE = RGBColor(26, 58, 92)
LIGHT_BLUE = RGBColor(240, 244, 249)
TEXT = RGBColor(34, 34, 34)
MUTED = RGBColor(96, 96, 96)
GREEN = RGBColor(27, 122, 64)


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert a Reveal-style HTML deck to PPTX.")
    parser.add_argument("html", type=Path)
    parser.add_argument("pptx", type=Path)
    args = parser.parse_args()
    convert(args.html, args.pptx)


def convert(html_path: Path, pptx_path: Path) -> None:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    sections = soup.select(".slides > section")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _set_background(slide, RGBColor(250, 251, 252))
        _add_footer(slide, idx, len(sections))

        if "title-slide" in section.get("class", []):
            _render_title_slide(slide, section)
            continue

        title = _first_heading(section)
        _add_title(slide, title or f"Slide {idx}")

        images = _image_paths(section, html_path.parent)
        tables = section.find_all("table")
        formulas = section.select(".formula-box")
        metric_cards = section.select(".metric-card")
        bullets = _extract_bullets(section)
        paragraphs = _extract_paragraphs(section)

        right_has_visual = bool(images)
        left_width = 7.15 if right_has_visual else 12.15
        y = 1.15

        if paragraphs:
            y = _add_paragraphs(slide, paragraphs[:3], Inches(0.65), Inches(y), Inches(left_width), max_lines=5)

        if bullets:
            y = _add_bullets(slide, bullets[:9], Inches(0.75), Inches(y + 0.08), Inches(left_width), Inches(4.9))

        if formulas:
            formula_text = "\n\n".join(_clean_text(f.get_text("\n")) for f in formulas)
            _add_formula_box(slide, formula_text, Inches(0.72), Inches(max(y + 0.12, 4.75)), Inches(left_width), Inches(1.2))

        if metric_cards:
            _add_metric_cards(slide, metric_cards[:4], Inches(0.75), Inches(1.35), Inches(left_width), Inches(3.2))

        if tables:
            _add_html_table(slide, tables[0], Inches(0.7), Inches(1.25), Inches(11.9), Inches(4.8))

        if images:
            if len(images) == 1:
                x = Inches(8.15) if left_width < 10 else Inches(2.1)
                w = Inches(4.55) if left_width < 10 else Inches(9.2)
                _add_image(slide, images[0], x, Inches(1.35), w, Inches(4.9))
            else:
                _add_image_grid(slide, images[:4], Inches(7.55), Inches(1.15), Inches(5.3), Inches(5.2))

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    print(f"Wrote {pptx_path} ({len(sections)} slides)")


def _render_title_slide(slide, section: Tag) -> None:
    _set_background(slide, RGBColor(245, 248, 252))
    h1 = section.find("h1")
    title = _clean_text(h1.get_text("\n")) if h1 else "Presentation"
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(1.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = BLUE

    subtitles = [_clean_text(p.get_text(" ")) for p in section.find_all("p")]
    if subtitles:
        sub = slide.shapes.add_textbox(Inches(1.4), Inches(3.75), Inches(10.5), Inches(1.0))
        stf = sub.text_frame
        stf.clear()
        for i, text in enumerate(subtitles):
            para = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = text
            run.font.size = Pt(19)
            run.font.color.rgb = MUTED

    line = slide.shapes.add_shape(1, Inches(3.3), Inches(3.45), Inches(6.7), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE


def _add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.25), Inches(0.62))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = BLUE
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(0.96), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE


def _add_footer(slide, idx: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.04), Inches(1.0), Inches(0.25))
    tf = box.text_frame
    tf.text = f"{idx}/{total}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].runs[0].font.size = Pt(9)
    tf.paragraphs[0].runs[0].font.color.rgb = MUTED


def _add_paragraphs(slide, paragraphs: list[str], x, y, w, max_lines: int) -> float:
    cleaned = [p for p in paragraphs if p]
    if not cleaned:
        return y.inches
    text = "\n\n".join(cleaned)
    box = slide.shapes.add_textbox(x, y, w, Inches(min(1.4, 0.32 * max_lines)))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.text = text
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT
    return y.inches + min(1.4, 0.32 * max_lines)


def _add_bullets(slide, bullets: list[tuple[int, str]], x, y, w, h) -> float:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (level, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = min(level, 2)
        p.font.size = Pt(15 if level == 0 else 12)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)
    return y.inches + h.inches


def _add_formula_box(slide, text: str, x, y, w, h) -> None:
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = RGBColor(184, 207, 232)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Courier New"
    run.font.size = Pt(12)
    run.font.color.rgb = BLUE


def _add_metric_cards(slide, cards: list[Tag], x, y, w, h) -> None:
    card_w = w / 2 - Inches(0.12)
    card_h = h / 2 - Inches(0.12)
    for i, card in enumerate(cards):
        cx = x + (card_w + Inches(0.24)) * (i % 2)
        cy = y + (card_h + Inches(0.24)) * (i // 2)
        shape = slide.shapes.add_shape(1, cx, cy, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.color.rgb = RGBColor(184, 207, 232)
        tf = shape.text_frame
        tf.clear()
        text = _clean_text(card.get_text("\n"))
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = BLUE


def _add_html_table(slide, table: Tag, x, y, w, h) -> None:
    rows = []
    for tr in table.find_all("tr"):
        cells = [_clean_text(c.get_text(" ")) for c in tr.find_all(["th", "td"])]
        if cells:
            rows.append(cells)
    if not rows:
        return
    cols = max(len(r) for r in rows)
    shape = slide.shapes.add_table(len(rows), cols, x, y, w, h)
    ppt_table = shape.table
    for r_idx, row in enumerate(rows):
        for c_idx in range(cols):
            cell = ppt_table.cell(r_idx, c_idx)
            cell.text = row[c_idx] if c_idx < len(row) else ""
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(255, 255, 255) if r_idx == 0 else TEXT
                run.font.bold = r_idx == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r_idx == 0 else RGBColor(244, 247, 251)


def _add_image_grid(slide, paths: list[Path], x, y, w, h) -> None:
    cell_w = w / 2 - Inches(0.08)
    cell_h = h / 2 - Inches(0.08)
    for i, path in enumerate(paths):
        _add_image(slide, path, x + (cell_w + Inches(0.16)) * (i % 2), y + (cell_h + Inches(0.16)) * (i // 2), cell_w, cell_h)


def _add_image(slide, path: Path, x, y, w, h) -> None:
    if not path.exists():
        return
    with Image.open(path) as image:
        aspect = image.width / image.height
    box_aspect = w / h
    if aspect >= box_aspect:
        final_w = w
        final_h = w / aspect
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * aspect
        final_x = x + (w - final_w) / 2
        final_y = y
    slide.shapes.add_picture(str(path), final_x, final_y, width=final_w, height=final_h)


def _first_heading(section: Tag) -> str | None:
    heading = section.find(["h1", "h2"])
    return _clean_text(heading.get_text("\n")) if heading else None


def _image_paths(section: Tag, base_dir: Path) -> list[Path]:
    paths = []
    for img in section.find_all("img"):
        src = img.get("src")
        if src and not src.startswith(("http://", "https://")):
            paths.append((base_dir / src).resolve())
    return paths


def _extract_paragraphs(section: Tag) -> list[str]:
    paragraphs = []
    for p in section.find_all("p"):
        if "caption" in p.get("class", []) or "subtitle" in p.get("class", []):
            continue
        if p.find_parent(["li", "table"]) or p.find_parent(class_="formula-box") or p.find_parent(class_="metric-card"):
            continue
        text = _clean_text(p.get_text(" "))
        if text:
            paragraphs.append(text)
    return paragraphs


def _extract_bullets(section: Tag) -> list[tuple[int, str]]:
    bullets: list[tuple[int, str]] = []
    for ul in section.find_all("ul"):
        if ul.find_parent("li"):
            continue
        for li in ul.find_all("li", recursive=False):
            _collect_li(li, 0, bullets)
    return bullets


def _collect_li(li: Tag, level: int, bullets: list[tuple[int, str]]) -> None:
    pieces = []
    for child in li.children:
        if isinstance(child, NavigableString):
            pieces.append(str(child))
        elif isinstance(child, Tag) and child.name != "ul":
            pieces.append(child.get_text(" "))
    text = _clean_text(" ".join(pieces))
    if text:
        bullets.append((level, text))
    for nested in li.find_all("ul", recursive=False):
        for nested_li in nested.find_all("li", recursive=False):
            _collect_li(nested_li, level + 1, bullets)


def _set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _clean_text(text: str) -> str:
    return " ".join(text.replace("\xa0", " ").split())


if __name__ == "__main__":
    main()
